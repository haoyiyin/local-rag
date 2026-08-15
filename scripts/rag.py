#!/usr/bin/env python3
"""local-rag: ingest, ask, and manage Chroma collections."""

import argparse, glob, hashlib, json, os, sys, textwrap, tempfile, time, urllib.request
from datetime import datetime, timezone
from pathlib import Path

import chromadb

CHROMA_PATH = os.environ.get("CHROMA_PATH", os.path.expanduser("~/.chroma/local-rag"))
COLL = os.environ.get("KB_COLLECTION", "default")

# Embedding model: default = embeddinggemma-300m Q8_0 GGUF (~334MB, true quantized).
# Set KB_EMBED_MODEL to a local dir or HuggingFace/ModelScope sentence-transformers
# model name to override with the PyTorch backend.
_Q8_GGUF_REPO = "ggml-org/embeddinggemma-300M-GGUF"
_Q8_GGUF_FILE = "embeddinggemma-300M-Q8_0.gguf"
_EMBED_FN = None
EMBED_MODEL = None


def _modelscope_cache_dir(repo):
    """Mirror the cache layout modelscope snapshot_download creates."""
    return os.path.expanduser(f"~/.cache/modelscope/models/{repo.replace('/', '--')}/snapshots/master")


def _modelscope_file_url(repo, filename):
    return f"https://modelscope.cn/models/{repo}/resolve/master/{filename}"


def _ensure_q8_gguf():
    """Download embeddinggemma-300M Q8_0 GGUF on first use if missing."""
    dest_dir = _modelscope_cache_dir(_Q8_GGUF_REPO)
    dest = os.path.join(dest_dir, _Q8_GGUF_FILE)
    if os.path.exists(dest):
        return dest

    os.makedirs(dest_dir, exist_ok=True)
    tmp = dest + ".part"
    url = _modelscope_file_url(_Q8_GGUF_REPO, _Q8_GGUF_FILE)
    print(f"Downloading embedding model {_Q8_GGUF_FILE} (~334MB)...", file=sys.stderr)
    req = urllib.request.Request(url, headers={"User-Agent": "local-rag/1.0"})
    with urllib.request.urlopen(req) as r, open(tmp, "wb") as out:
        total = int(r.headers.get("Content-Length") or 0)
        done = 0
        while True:
            buf = r.read(1024 * 256)
            if not buf:
                break
            out.write(buf)
            done += len(buf)
            if total:
                pct = done * 100 // total
                print(f"\r  {pct:3d}% ({done // (1024*1024)}/{total // (1024*1024)} MB)",
                      file=sys.stderr, end="", flush=True)
        if total:
            print(file=sys.stderr)
    if total and os.path.getsize(tmp) < total:
        os.unlink(tmp)
        sys.exit("Download incomplete. Re-run to resume.")
    os.replace(tmp, dest)
    return dest


def _build_embed_fn():
    """Lazy: return (embedding_function, model_label)."""
    override = os.environ.get("KB_EMBED_MODEL", "").strip()
    if override:
        try:
            from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
        except ImportError:
            sys.exit("KB_EMBED_MODEL requires sentence-transformers. Run: pip install sentence-transformers")
        return SentenceTransformerEmbeddingFunction(model_name=override), override

    model_path = _ensure_q8_gguf()
    try:
        from llama_cpp import Llama
    except ImportError:
        sys.exit("Missing llama-cpp-python. Run: pip install llama-cpp-python")

    class _LlamaCppEmbeddingFunction:
        """GGUF Q8 backend. Llama.cpp applies EmbeddingGemma's ST Dense head,
        so output matches the original 768-dim normalized embeddings."""
        def __init__(self, path):
            self._path = path
            self._llm = None

        def _load(self):
            if self._llm is None:
                self._llm = Llama(model_path=self._path, embedding=True, verbose=False)
            return self._llm

        def __call__(self, docs):
            llm = self._load()
            out = []
            for d in docs:
                emb = llm.embed(d, normalize=True)[0]
                out.append(list(emb))
            return out

    return _LlamaCppEmbeddingFunction(model_path), f"gguf:{model_path}"


def _embedding():
    """Cached embedding function. Model downloads/loads on first use."""
    global _EMBED_FN, EMBED_MODEL
    if _EMBED_FN is None:
        _EMBED_FN, EMBED_MODEL = _build_embed_fn()
    return _EMBED_FN


def _client():
    return chromadb.PersistentClient(path=CHROMA_PATH)

def _coll_required(c, name):
    """Get collection by name, fail if doesn't exist. Prevents silent typo errors."""
    if not name:
        sys.exit("error: collection name required (use -c <name>)")
    try:
        return c.get_collection(name, embedding_function=_embedding())
    except Exception:
        existing = [col.name for col in c.list_collections()]
        sys.exit(f"error: collection '{name}' does not exist. Available: {existing}")


def _coll_required_or_auto(c, name):
    """Get collection — fail strict, or fall back to default."""
    if name:
        return _coll_required(c, name)
    return c.get_or_create_collection(COLL, embedding_function=_embedding())

def _id(source, idx):
    return hashlib.sha1(f"{source}\0{idx}".encode()).hexdigest()[:16]

def _chunks(text, max_chars=1200):
    """Split text on blank lines, cap each chunk."""
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    cur = []
    cur_len = 0
    for p in paragraphs:
        if cur_len + len(p) + 1 > max_chars and cur:
            chunks.append("\n\n".join(cur))
            cur, cur_len = [], 0
        cur.append(p)
        cur_len += len(p) + 1
    if cur:
        chunks.append("\n\n".join(cur))
    return chunks or [text[:max_chars]]


ARCHIVE_EXTS = {".zip", ".tar", ".gz", ".tgz", ".bz2", ".tbz2", ".xz", ".txz", ".7z", ".rar"}

def _is_archive(path):
    """Check if file is a supported archive format."""
    lower = path.lower()
    for ext in [".tar.gz", ".tar.bz2", ".tar.xz", ".tgz", ".tbz2", ".txz"]:
        if lower.endswith(ext):
            return True
    return Path(path).suffix.lower() in ARCHIVE_EXTS

def _extract_archive(path, dest):
    """Extract archive to dest dir. Returns list of extracted file paths."""
    lower = path.lower()
    if lower.endswith(".zip"):
        import zipfile
        with zipfile.ZipFile(path, "r") as zf:
            zf.extractall(dest)
    elif lower.endswith((".tar.gz", ".tgz")):
        import tarfile
        with tarfile.open(path, "r:gz") as tf:
            tf.extractall(dest, filter="data")
    elif lower.endswith((".tar.bz2", ".tbz2")):
        import tarfile
        with tarfile.open(path, "r:bz2") as tf:
            tf.extractall(dest, filter="data")
    elif lower.endswith((".tar.xz", ".txz")):
        import tarfile
        with tarfile.open(path, "r:xz") as tf:
            tf.extractall(dest, filter="data")
    elif lower.endswith(".tar"):
        import tarfile
        with tarfile.open(path, "r") as tf:
            tf.extractall(dest, filter="data")
    elif lower.endswith(".7z"):
        import py7zr
        with py7zr.SevenZipFile(path, "r") as sz:
            sz.extractall(dest)
    elif lower.endswith(".rar"):
        import rarfile
        with rarfile.RarFile(path) as rf:
            rf.extractall(dest)
    else:
        sys.exit(f"Unsupported archive format: {path}")
    # Collect all non-hidden files
    files = []
    for root, dirs, fnames in os.walk(dest):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for f in fnames:
            if not f.startswith("."):
                files.append(os.path.join(root, f))
    return sorted(files)


def _read_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            sys.exit("PDF support requires pypdf. Run: pip install pypdf")
        reader = PdfReader(path)
        text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        if len(text.strip()) > 50:
            return text
        try:
            import pymupdf
            import pytesseract
            from PIL import Image
            import io
        except ImportError:
            print("WARNING: image PDF detected but pymupdf/pytesseract not available; returning partial text", file=sys.stderr)
            return text
        doc = pymupdf.open(path)
        ocr_parts = []
        for page in doc:
            mat = pymupdf.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            page_text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            if page_text.strip():
                ocr_parts.append(page_text.strip())
        doc.close()
        return "\n\n".join(ocr_parts) or text
    elif ext in (".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac", ".wma", ".opus"):
        try:
            import whisper
        except ImportError:
            sys.exit("Audio support requires openai-whisper. Run: pip install openai-whisper")
        model = whisper.load_model("base")
        result = model.transcribe(path, fp16=False)
        lang = result.get("language", "")
        text = result.get("text", "").strip()
        if text:
            return f"[audio lang={lang}]\n{text}"
        return ""
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"):
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            sys.exit("Image OCR requires pytesseract + Pillow. Run: pip install pytesseract Pillow")
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng").strip()
        return f"[image ocr]\n{text}" if text else ""
    elif ext == ".docx":
        try:
            import docx
        except ImportError:
            sys.exit("DOCX support requires python-docx. Run: pip install python-docx")
        doc = docx.Document(path)
        return "\n\n".join(p.text for p in doc.paragraphs)
    elif ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            sys.exit("PPTX support requires python-pptx. Run: pip install python-pptx")
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                parts.append("\n".join(texts))
        return "\n\n".join(parts)
    elif ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
        except ImportError:
            sys.exit("XLSX support requires openpyxl. Run: pip install openpyxl")
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    rows.append("\t".join(vals))
            if rows:
                parts.append(f"[{sheet}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    elif ext == ".csv":
        import csv, io
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = ["\t".join(row) for row in reader if any(c.strip() for c in row)]
        return "\n".join(rows)
    elif ext in (".html", ".htm"):
        from html.parser import HTMLParser
        class _HTMLText(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    self._parts.append(data)
            def text(self):
                return " ".join(self._parts)
        with open(path, encoding="utf-8", errors="replace") as f:
            parser = _HTMLText()
            parser.feed(f.read())
        return parser.text()
    elif ext in (".json", ".jsonl"):
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)
    elif ext in (".yaml", ".yml"):
        try:
            import yaml
        except ImportError:
            sys.exit("YAML support requires PyYAML. Run: pip install pyyaml")
        with open(path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
        return json.dumps(data, ensure_ascii=False, indent=2) if not isinstance(data, str) else data
    elif ext == ".odt":
        try:
            from odf.opendocument import load
            from odf.text import P
        except ImportError:
            sys.exit("ODT support requires odfpy. Run: pip install odfpy")
        doc = load(path)
        parts = [elem.firstChild.data if elem.firstChild else "" for elem in doc.getElementsByType(P)]
        return "\n\n".join(p for p in parts if p.strip())
    elif ext == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from html.parser import HTMLParser
        except ImportError:
            sys.exit("EPUB support requires ebooklib. Run: pip install ebooklib")
        class _EpubText(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    self._parts.append(data)
            def text(self):
                return " ".join(self._parts)
        book = epub.read_epub(path)
        parts = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                p = _EpubText()
                p.feed(item.get_content().decode("utf-8", errors="replace"))
                t = p.text().strip()
                if t:
                    parts.append(t)
        return "\n\n".join(parts)
    elif ext == ".rtf":
        from striprtf.striprtf import rtf_to_text
        with open(path, encoding="utf-8", errors="replace") as f:
            return rtf_to_text(f.read())
    elif ext == ".msg":
        try:
            import extract_msg
        except ImportError:
             sys.exit("MSG support requires extract-msg. Run: pip install extract-msg")
        msg = extract_msg.Message(path)
        parts = [f"From: {msg.sender}", f"To: {msg.to}", f"Subject: {msg.subject}", f"Date: {msg.date}"]
        body = msg.body or ""
        parts.append(body)
        return "\n\n".join(p for p in parts if p)
    elif ext == ".ipynb":
        with open(path, encoding="utf-8") as f:
            nb = json.load(f)
        parts = []
        for cell in nb.get("cells", []):
            ctype = cell.get("cell_type", "")
            src = "".join(cell.get("source", []))
            if ctype == "markdown":
                parts.append(src)
            elif ctype == "code":
                outs = []
                for out in cell.get("outputs", []):
                    if "text" in out:
                        outs.append("".join(out["text"]))
                    elif "data" in out:
                        for k, v in out["data"].items():
                            if k.startswith("text/"):
                                outs.append("".join(v) if isinstance(v, list) else v)
                if outs:
                    parts.append(f"[code]\n{src}\n[output]\n" + "\n".join(outs))
                else:
                    parts.append(f"[code]\n{src}")
        return "\n\n".join(parts)
    else:
        with open(path) as f:
            return f.read()


# ─── Collection management commands ───

def cmd_list(args):
    """List all collections with record counts."""
    c = _client()
    cols = c.list_collections(limit=args.limit, offset=args.offset)
    if not cols:
        print("no collections")
        return
    rows = []
    for col in cols:
        n = col.count()
        meta = col.metadata or {}
        rows.append({"name": col.name, "id": str(col.id), "count": n, "metadata": meta})
    print(json.dumps(rows, ensure_ascii=False, indent=2))


def cmd_create(args):
    """Create a new collection with default embedding function."""
    c = _client()
    metadata = {}
    if args.meta:
        for kv in args.meta:
            k, _, v = kv.partition("=")
            # try to parse as JSON value, else keep as string
            try:
                metadata[k] = json.loads(v)
            except (json.JSONDecodeError, TypeError):
                metadata[k] = v
    col = c.create_collection(
        args.name,
        embedding_function=_embedding(),
        metadata=metadata or None,
    )
    print(json.dumps({"created": col.name, "id": str(col.id), "embedding_model": EMBED_MODEL}))


def cmd_delete(args):
    """Delete a collection by name."""
    c = _client()
    c.delete_collection(args.name)
    print(json.dumps({"deleted": args.name}))


def cmd_info(args):
    """Show collection details: count, metadata, sample records."""
    c = _client()
    col = c.get_collection(args.name)
    n = col.count()
    peek = col.peek(limit=args.sample)
    p_ids = peek["ids"] if isinstance(peek, dict) else peek.ids
    p_docs = peek["documents"] if isinstance(peek, dict) else peek.documents
    p_metas = peek["metadatas"] if isinstance(peek, dict) else peek.metadatas
    info = {
        "name": col.name,
        "id": str(col.id),
        "count": n,
        "metadata": col.metadata,
        "sample_ids": p_ids,
        "sample_documents": p_docs,
        "sample_metadatas": p_metas,
    }
    print(json.dumps(info, ensure_ascii=False, indent=2))


def cmd_rename(args):
    """Rename a collection."""
    c = _client()
    col = c.get_collection(args.old_name)
    col.modify(name=args.new_name)
    print(json.dumps({"renamed": f"{args.old_name} → {args.new_name}"}))


def cmd_purge(args):
    """Delete all records in a collection (keeps the collection)."""
    c = _client()
    col = c.get_collection(args.name)
    n = col.count()
    if n == 0:
        print(json.dumps({"purged": args.name, "deleted": 0}))
        return
    all_ids = []
    while True:
        batch = col.get(limit=1000, include=[])
        ids = batch["ids"] if isinstance(batch, dict) else batch.ids
        if not ids:
            break
        all_ids.extend(ids)
        if len(ids) < 1000:
            break
    col.delete(ids=all_ids)
    print(json.dumps({"purged": args.name, "deleted": len(all_ids)}))


# ─── Data commands ───

def cmd_ingest(args):
    c = _client()
    col = _coll_required(c, args.collection)

    # Archive: extract and ingest all files inside
    if args.path != "-" and _is_archive(args.path):
        archive_name = os.path.basename(args.path)
        tag = args.tag or archive_name
        with tempfile.TemporaryDirectory(prefix="localrag_") as tmpdir:
            files = _extract_archive(args.path, tmpdir)
            if not files:
                print(json.dumps({"error": "archive empty", "archive": archive_name}))
                return
            total_chunks = 0
            ingested = []
            for fpath in files:
                ext = os.path.splitext(fpath)[1].lower()
                if ext in (".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go",
                           ".rs", ".rb", ".php", ".swift", ".kt", ".sh", ".bat",
                           ".exe", ".dll", ".so", ".dylib", ".bin", ".png", ".jpg",
                           ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp",
                           ".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac",
                           ".wma", ".opus", ".mp4", ".avi", ".mov", ".mkv"):
                    continue  # skip binary/code files
                try:
                    text = _read_file(fpath)
                except Exception:
                    continue
                if not text or not text.strip():
                    continue
                source = os.path.relpath(fpath, tmpdir)
                chunks = _chunks(text)
                ids = [_id(f"{archive_name}/{source}", i) for i in range(len(chunks))]
                metas = [
                    {"source": f"{archive_name}/{source}", "tag": tag,
                     "chunk_idx": i, "archive": archive_name,
                     "ingested_at": datetime.now(timezone.utc).isoformat()}
                    for i in range(len(chunks))
                ]
                try:
                    col.delete(where={"source": f"{archive_name}/{source}"})
                except Exception:
                    pass
                col.add(documents=chunks, metadatas=metas, ids=ids)
                total_chunks += len(chunks)
                ingested.append({"file": source, "chunks": len(chunks)})
            print(json.dumps({
                "archive": archive_name,
                "files_ingested": len(ingested),
                "total_chunks": total_chunks,
                "details": ingested,
                "collection": args.collection or COLL,
            }, ensure_ascii=False))
        return

    # Single file or stdin
    source = args.path
    if args.path == "-":
        text = sys.stdin.read()
        source = args.source or "stdin"
    else:
        text = _read_file(args.path)
        source = args.source or os.path.basename(args.path)

    chunks = _chunks(text)
    ids = [_id(source, i) for i in range(len(chunks))]
    metas = [
        {"source": source, "tag": args.tag or "", "chunk_idx": i, "ingested_at": datetime.now(timezone.utc).isoformat()}
        for i in range(len(chunks))
    ]

    try:
        col.delete(where={"source": source})
    except Exception:
        pass

    col.add(documents=chunks, metadatas=metas, ids=ids)
    print(json.dumps({"chunks_added": len(chunks), "source": source, "collection": args.collection or COLL}))


def cmd_ingest_text(args):
    c = _client()
    col = _coll_required(c, args.collection)

    text = args.text
    source = args.source or "inline"
    chunks = _chunks(text)
    ids = [_id(source, i) for i in range(len(chunks))]
    metas = [
        {"source": source, "tag": args.tag or "", "chunk_idx": i, "ingested_at": datetime.now(timezone.utc).isoformat()}
        for i in range(len(chunks))
    ]

    try:
        col.delete(where={"source": source})
    except Exception:
        pass

    col.add(documents=chunks, metadatas=metas, ids=ids)
    print(json.dumps({"chunks_added": len(chunks), "source": source, "collection": args.collection}))


def cmd_ask(args):
    c = _client()
    col = _coll_required(c, args.collection)

    where = None
    if args.tag:
        where = {"tag": args.tag}

    results = col.query(
        query_texts=[args.query],
        n_results=args.k,
        where=where,
        include=["documents", "metadatas", "distances"],
    )

    out = []
    if results["documents"] and results["documents"][0]:
        for i, doc in enumerate(results["documents"][0]):
            meta = results["metadatas"][0][i] if results["metadatas"] else {}
            out.append({
                "source": meta.get("source", "?"),
                "tag": meta.get("tag", ""),
                "chunk": doc,
                "score": round(results["distances"][0][i], 4) if results["distances"] else 0.0,
            })
    print(json.dumps(out, ensure_ascii=False))


def cmd_smoke(args):
    c = _client()
    col = _coll_required_or_auto(c, None)
    col.upsert(documents=["smoke test"], ids=["__smoke"], metadatas=[{"source": "__smoke"}])
    r = col.query(query_texts=["smoke"], n_results=1)
    hits = len(r["documents"][0]) if r["documents"] else 0
    col.delete(ids=["__smoke"])
    print(json.dumps({"smoke_ok": True, "hits": hits}))


def main():
    p = argparse.ArgumentParser(description="local-rag: Chroma knowledge base CLI")
    sub = p.add_subparsers(dest="cmd", required=True)

    # --- Collection management ---
    pl = sub.add_parser("list", help="List all collections")
    pl.add_argument("--limit", type=int, default=100, help="Max collections to list")
    pl.add_argument("--offset", type=int, default=0, help="Offset for pagination")

    pc = sub.add_parser("create", help="Create a new collection")
    pc.add_argument("name", help="Collection name")
    pc.add_argument("--meta", nargs="*", help="Metadata as key=value pairs")

    pd = sub.add_parser("delete", help="Delete a collection")
    pd.add_argument("name", help="Collection name to delete")

    pi_info = sub.add_parser("info", help="Show collection details")
    pi_info.add_argument("name", help="Collection name")
    pi_info.add_argument("--sample", type=int, default=3, help="Number of sample records")

    pr = sub.add_parser("rename", help="Rename a collection")
    pr.add_argument("old_name", help="Current name")
    pr.add_argument("new_name", help="New name")

    pp = sub.add_parser("purge", help="Delete all records in a collection")
    pp.add_argument("name", help="Collection name")

    # --- Data commands (require -c, refuse if collection doesn't exist) ---
    pi = sub.add_parser("ingest", help="Store a file (or stdin with -)")
    pi.add_argument("path", help="File path or '-' for stdin")
    pi.add_argument("--tag", default="", help="Optional tag for filtering")
    pi.add_argument("--source", default=None, help="Override source label")
    pi.add_argument("--collection", "-c", required=True, help="Target collection (must already exist — use 'list' to see available)")

    pit = sub.add_parser("ingest-text", help="Store explicit text string")
    pit.add_argument("text", help="Text to store")
    pit.add_argument("--tag", default="", help="Optional tag")
    pit.add_argument("--source", default=None, help="Source label")
    pit.add_argument("--collection", "-c", required=True, help="Target collection (must already exist)")

    pa = sub.add_parser("ask", help="Retrieve relevant chunks")
    pa.add_argument("query", help="Search query")
    pa.add_argument("--k", type=int, default=5, help="Number of results (default 5)")
    pa.add_argument("--tag", default=None, help="Filter by tag")
    pa.add_argument("--collection", "-c", required=True, help="Target collection (must already exist)")

    ps = sub.add_parser("smoke", help="Quick connectivity test")

    args = p.parse_args()

    cmds = {
        "list": cmd_list, "create": cmd_create, "delete": cmd_delete,
        "info": cmd_info, "rename": cmd_rename, "purge": cmd_purge,
        "ingest": cmd_ingest, "ingest-text": cmd_ingest_text,
        "ask": cmd_ask, "smoke": cmd_smoke,
    }
    try:
        cmds[args.cmd](args)
    except Exception as e:
        msg = str(e)
        if "No module named" in msg:
            sys.exit(f"Missing dependency: {msg}. Run: pip install chromadb")
        raise


if __name__ == "__main__":
    main()
